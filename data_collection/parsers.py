from __future__ import annotations

from pathlib import Path
from zipfile import ZipFile

from common.file_ops import read_text


def _fallback(path: Path, reason: str) -> str:
    return f"# {path.stem}\n\n> 解析降级：{reason}\n\n原始文件：{path.name}\n"


def parse_markdown(path: Path) -> str:
    return read_text(path)


def parse_text(path: Path) -> str:
    return f"# {path.stem}\n\n" + read_text(path)


def parse_pdf(path: Path) -> str:
    try:
        from pypdf import PdfReader  # type: ignore

        reader = PdfReader(str(path))
        pages = []
        for idx, page in enumerate(reader.pages, start=1):
            pages.append(f"## Page {idx}\n\n{page.extract_text() or ''}")
        return f"# {path.stem}\n\n" + "\n\n".join(pages)
    except Exception as exc:
        return _fallback(path, f"PDF 文本解析不可用或失败（{exc.__class__.__name__}）")


def parse_docx(path: Path) -> str:
    try:
        import docx  # type: ignore

        doc = docx.Document(str(path))
        lines = [f"# {path.stem}"]
        for paragraph in doc.paragraphs:
            text = paragraph.text.strip()
            if text:
                lines.append(text)
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                lines.append("| " + " | ".join(cells) + " |")
        return "\n\n".join(lines) + "\n"
    except Exception as exc:
        if path.suffix.lower() == ".docx":
            try:
                with ZipFile(path) as zf:
                    xml = zf.read("word/document.xml").decode("utf-8", errors="ignore")
                import re

                text = re.sub(r"<[^>]+>", " ", xml)
                return f"# {path.stem}\n\n{text}\n"
            except Exception:
                pass
        return _fallback(path, f"Word 文档解析不可用或失败（{exc.__class__.__name__}）")


def parse_excel(path: Path) -> str:
    try:
        import pandas as pd  # type: ignore

        sheets = pd.read_excel(path, sheet_name=None)
        parts = [f"# {path.stem}"]
        for name, frame in sheets.items():
            parts.append(f"## {name}\n\n" + frame.to_markdown(index=False))
        return "\n\n".join(parts) + "\n"
    except Exception as exc:
        return _fallback(path, f"Excel 表格解析不可用或失败（{exc.__class__.__name__}）")


def parse_ppt(path: Path) -> str:
    try:
        from pptx import Presentation  # type: ignore

        prs = Presentation(str(path))
        parts = [f"# {path.stem}"]
        for idx, slide in enumerate(prs.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            parts.append(f"## Slide {idx}\n\n" + "\n\n".join(texts))
        return "\n\n".join(parts) + "\n"
    except Exception as exc:
        return _fallback(path, f"PPT 演示文稿解析不可用或失败（{exc.__class__.__name__}）")


def parse_image(path: Path) -> str:
    try:
        import pytesseract  # type: ignore
        from PIL import Image  # type: ignore

        text = pytesseract.image_to_string(Image.open(path))
        return f"# {path.stem}\n\n![{path.stem}]({path})\n\n{text}\n"
    except Exception as exc:
        return _fallback(path, f"图片 OCR 不可用或失败（{exc.__class__.__name__}）")


def parse_to_markdown(path: str | Path) -> str:
    p = Path(path)
    suffix = p.suffix.lower()
    if suffix in {".md", ".markdown"}:
        return parse_markdown(p)
    if suffix in {".txt", ".log", ".csv"}:
        return parse_text(p)
    if suffix == ".pdf":
        return parse_pdf(p)
    if suffix in {".doc", ".docx"}:
        return parse_docx(p)
    if suffix in {".xls", ".xlsx"}:
        return parse_excel(p)
    if suffix in {".ppt", ".pptx"}:
        return parse_ppt(p)
    if suffix in {".jpg", ".jpeg", ".png", ".webp", ".bmp", ".tiff"}:
        return parse_image(p)
    return _fallback(p, f"暂不支持的文件类型 {suffix or 'unknown'}")
